import logging
import asyncio
from typing import List, Dict, Any, Optional, Tuple, Union
import base64
import io
from pathlib import Path
import tempfile
import uuid

import pandas as pd
import numpy as np
from PIL import Image
import pytesseract
import fitz  # PyMuPDF
from langchain.schema import Document
import openpyxl
from pptx import Presentation
import cv2

from app.core.neo4j_client import Neo4jClient
from app.core.exceptions import DocumentProcessingError
from app.config.settings import get_settings
from app.utils.llm_clients import LLMClientFactory

logger = logging.getLogger(__name__)

class MultiModalProcessor:
    """Advanced multi-modal document processing"""
    
    def __init__(self, neo4j_client: Neo4jClient):
        self.neo4j = neo4j_client
        self.settings = get_settings()
        self.llm_factory = LLMClientFactory()
        
        # Supported file types and their processors
        self.processors = {
            '.pdf': self._process_pdf_advanced,
            '.xlsx': self._process_excel,
            '.xls': self._process_excel,
            '.pptx': self._process_powerpoint,
            '.ppt': self._process_powerpoint,
            '.docx': self._process_word,
            '.doc': self._process_word,
            '.csv': self._process_csv,
            '.json': self._process_json,
            '.html': self._process_html,
            '.xml': self._process_xml,
            '.png': self._process_image,
            '.jpg': self._process_image,
            '.jpeg': self._process_image,
            '.tiff': self._process_image,
            '.bmp': self._process_image
        }
    
    async def process_multimodal_document(self, file_path: str, document_id: str) -> Dict[str, Any]:
        """Process multi-modal document with advanced extraction"""
        try:
            file_extension = Path(file_path).suffix.lower()
            
            if file_extension not in self.processors:
                raise DocumentProcessingError(f"Unsupported file type: {file_extension}")
            
            # Get processor function
            processor = self.processors[file_extension]
            
            # Process document
            processing_result = await processor(file_path, document_id)
            
            # Extract metadata
            metadata = await self._extract_file_metadata(file_path)
            processing_result["file_metadata"] = metadata
            
            # Analyze document structure
            structure_analysis = await self._analyze_document_structure(processing_result)
            processing_result["structure_analysis"] = structure_analysis
            
            return processing_result
            
        except Exception as e:
            logger.error(f"Multi-modal processing failed for {file_path}: {e}")
            raise DocumentProcessingError(f"Multi-modal processing failed: {e}")
    
    async def _process_pdf_advanced(self, file_path: str, document_id: str) -> Dict[str, Any]:
        """Advanced PDF processing with OCR, table extraction, and image analysis"""
        try:
            doc = fitz.open(file_path)
            result = {
                "text_content": [],
                "tables": [],
                "images": [],
                "metadata": {},
                "structure": {
                    "total_pages": len(doc),
                    "has_images": False,
                    "has_tables": False,
                    "text_density": 0.0
                }
            }
            
            total_text_length = 0
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                page_result = await self._process_pdf_page(page, page_num, document_id)
                
                # Accumulate results
                if page_result["text"]:
                    result["text_content"].append({
                        "page": page_num + 1,
                        "text": page_result["text"],
                        "extraction_method": page_result.get("method", "standard")
                    })
                    total_text_length += len(page_result["text"])
                
                if page_result["tables"]:
                    result["tables"].extend(page_result["tables"])
                    result["structure"]["has_tables"] = True
                
                if page_result["images"]:
                    result["images"].extend(page_result["images"])
                    result["structure"]["has_images"] = True
            
            # Calculate text density
            result["structure"]["text_density"] = total_text_length / max(len(doc), 1)
            
            # Extract document metadata
            result["metadata"] = {
                "title": doc.metadata.get("title", ""),
                "author": doc.metadata.get("author", ""),
                "subject": doc.metadata.get("subject", ""),
                "creator": doc.metadata.get("creator", ""),
                "creation_date": doc.metadata.get("creationDate", ""),
                "modification_date": doc.metadata.get("modDate", "")
            }
            
            doc.close()
            return result
            
        except Exception as e:
            logger.error(f"PDF processing failed: {e}")
            raise DocumentProcessingError(f"PDF processing failed: {e}")
    
    async def _process_pdf_page(self, page, page_num: int, document_id: str) -> Dict[str, Any]:
        """Process individual PDF page with OCR fallback"""
        page_result = {
            "text": "",
            "tables": [],
            "images": [],
            "method": "standard"
        }
        
        try:
            # Extract text using standard method
            text = page.get_text()
            
            # If text is sparse, use OCR
            if len(text.strip()) < 50:
                page_result["text"] = await self._ocr_pdf_page(page, page_num)
                page_result["method"] = "ocr"
            else:
                page_result["text"] = text
            
            # Extract tables
            tables = await self._extract_tables_from_page(page, page_num)
            page_result["tables"] = tables
            
            # Extract images
            images = await self._extract_images_from_page(page, page_num, document_id)
            page_result["images"] = images
            
            return page_result
            
        except Exception as e:
            logger.warning(f"Page {page_num} processing failed: {e}")
            return page_result
    
    async def _ocr_pdf_page(self, page, page_num: int) -> str:
        """Perform OCR on PDF page"""
        try:
            # Convert page to image
            mat = fitz.Matrix(2, 2)  # 2x zoom
            pix = page.get_pixmap(matrix=mat)
            img_data = pix.tobytes("png")
            
            # Convert to PIL Image
            image = Image.open(io.BytesIO(img_data))
            
            # Perform OCR
            ocr_text = await asyncio.to_thread(pytesseract.image_to_string, image, lang='eng')
            
            return ocr_text.strip()
            
        except Exception as e:
            logger.warning(f"OCR failed for page {page_num}: {e}")
            return ""
    
    async def _extract_tables_from_page(self, page, page_num: int) -> List[Dict[str, Any]]:
        """Extract tables from PDF page"""
        try:
            # Use pymupdf table extraction
            tables = page.find_tables()
            extracted_tables = []
            
            for i, table in enumerate(tables):
                try:
                    table_data = table.extract()
                    if table_data:
                        # Convert to structured format
                        df = pd.DataFrame(table_data[1:], columns=table_data[0] if table_data else [])
                        
                        extracted_tables.append({
                            "table_id": f"page_{page_num}_table_{i}",
                            "page": page_num + 1,
                            "data": df.to_dict('records'),
                            "columns": list(df.columns),
                            "rows": len(df),
                            "bbox": table.bbox  # Bounding box coordinates
                        })
                        
                except Exception as e:
                    logger.warning(f"Table {i} extraction failed on page {page_num}: {e}")
            
            return extracted_tables
            
        except Exception as e:
            logger.warning(f"Table extraction failed for page {page_num}: {e}")
            return []
    
    async def _extract_images_from_page(self, page, page_num: int, document_id: str) -> List[Dict[str, Any]]:
        """Extract and analyze images from PDF page"""
        try:
            image_list = page.get_images()
            extracted_images = []
            
            for i, img in enumerate(image_list):
                try:
                    # Get image data
                    xref = img[0]
                    base_image = page.parent.extract_image(xref)
                    image_bytes = base_image["image"]
                    
                    # Convert to PIL Image
                    image = Image.open(io.BytesIO(image_bytes))
                    
                    # Analyze image
                    image_analysis = await self._analyze_image(image, f"{document_id}_page_{page_num}_img_{i}")

                    extracted_images.append({
                        "image_id": f"page_{page_num}_image_{i}",
                        "page": page_num + 1,
                        "format": base_image["ext"],
                        "size": {"width": image.width, "height": image.height},
                        "analysis": image_analysis,
                        "has_text": len(image_analysis.get("ocr_text", "")) > 10
                    })
                    
                except Exception as e:
                    logger.warning(f"Image {i} extraction failed on page {page_num}: {e}")
            
            return extracted_images
            
        except Exception as e:
            logger.warning(f"Image extraction failed for page {page_num}: {e}")
            return []
    
    async def _analyze_image(self, image: Image.Image, image_id: str) -> Dict[str, Any]:
        """Analyze image content using OCR and vision models"""
        try:
            analysis = {}
            
            # OCR text extraction
            try:
                ocr_text = await asyncio.to_thread(pytesseract.image_to_string, image)
                analysis["ocr_text"] = ocr_text.strip()
            except Exception as e:
                logger.warning(f"OCR failed for image {image_id}: {e}")
                analysis["ocr_text"] = ""
            
            # Basic image properties
            analysis["properties"] = {
                "format": image.format,
                "mode": image.mode,
                "size": image.size,
                "has_transparency": image.mode in ("RGBA", "LA")
            }
            
            # Detect if image contains charts/diagrams
            analysis["content_type"] = await self._classify_image_content(image)
            
            # Extract text using vision model if available
            if self.settings.openai_api_key:
                analysis["vision_description"] = await self._describe_image_with_vision_model(image)
            
            return analysis
            
        except Exception as e:
            logger.warning(f"Image analysis failed for {image_id}: {e}")
            return {"error": str(e)}
    
    async def _classify_image_content(self, image: Image.Image) -> str:
        """Classify image content type (chart, diagram, photo, etc.)"""
        try:
            # Convert to numpy array for analysis
            img_array = np.array(image.convert('RGB'))
            
            # Simple heuristics for content classification
            # Check for high contrast (charts/diagrams tend to have high contrast)
            gray = cv2.cvtColor(img_array, cv2.COLOR_RGB2GRAY)
            contrast = gray.std()
            
            # Check for geometric shapes (lines, rectangles)
            edges = cv2.Canny(gray, 50, 150)
            lines = cv2.HoughLinesP(edges, 1, np.pi/180, threshold=50, minLineLength=30, maxLineGap=10)
            
            # Classify based on features
            if contrast > 50 and lines is not None and len(lines) > 10:
                return "diagram_chart"
            elif contrast > 30:
                return "text_document"
            elif contrast < 20:
                return "photograph"
            else:
                return "mixed_content"
                
        except Exception as e:
            logger.warning(f"Image classification failed: {e}")
            return "unknown"
    
    async def _describe_image_with_vision_model(self, image: Image.Image) -> str:
        """Use vision model to describe image content"""
        try:
            # Convert image to base64 for API
            buffered = io.BytesIO()
            image.save(buffered, format="PNG")
            img_base64 = base64.b64encode(buffered.getvalue()).decode()
            
            # Use OpenAI Vision API (GPT-4V)
            import openai
            client = openai.OpenAI(api_key=self.settings.openai_api_key)
            
            response = await asyncio.to_thread(
                client.chat.completions.create,
                model="gpt-4-vision-preview",
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": "Describe this image in detail, focusing on any text, charts, diagrams, or data visualizations present."
                            },
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/png;base64,{img_base64}"
                                }
                            }
                        ]
                    }
                ],
                max_tokens=500
            )
            
            return response.choices[0].message.content
            
        except Exception as e:
            logger.warning(f"Vision model description failed: {e}")
            return ""
    
    async def _process_excel(self, file_path: str, document_id: str) -> Dict[str, Any]:
        """Process Excel files with sheet and chart analysis"""
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            result = {
                "sheets": [],
                "charts": [],
                "pivot_tables": [],
                "metadata": {
                    "total_sheets": len(workbook.worksheets),
                    "sheet_names": [ws.title for ws in workbook.worksheets]
                }
            }
            
            for worksheet in workbook.worksheets:
                sheet_data = await self._process_excel_sheet(worksheet, document_id)
                result["sheets"].append(sheet_data)
                
                # Extract charts from sheet
                if hasattr(worksheet, '_charts'):
                    for chart in worksheet._charts:
                        chart_analysis = await self._analyze_excel_chart(chart, worksheet.title)
                        result["charts"].append(chart_analysis)
                
                # Detect pivot tables
                pivot_tables = await self._detect_pivot_tables(worksheet)
                result["pivot_tables"].extend(pivot_tables)
            
            workbook.close()
            return result
            
        except Exception as e:
            logger.error(f"Excel processing failed: {e}")
            raise DocumentProcessingError(f"Excel processing failed: {e}")
    
    async def _process_excel_sheet(self, worksheet, document_id: str) -> Dict[str, Any]:
        """Process individual Excel worksheet"""
        try:
            # Convert to DataFrame for easier processing
            data = []
            for row in worksheet.iter_rows(values_only=True):
                data.append(row)
            
            if not data:
                return {"name": worksheet.title, "empty": True}
            
            df = pd.DataFrame(data)
            
            # Clean and analyze data
            df = df.dropna(how='all').dropna(axis=1, how='all')  # Remove empty rows/columns
            
            # Detect header row
            header_row = await self._detect_header_row(df)
            
            if header_row is not None:
                df.columns = df.iloc[header_row]
                df = df.iloc[header_row + 1:]
                df = df.reset_index(drop=True)
            
            # Analyze data types and patterns
            analysis = await self._analyze_tabular_data(df, worksheet.title)
            
            return {
                "name": worksheet.title,
                "dimensions": {"rows": len(df), "columns": len(df.columns)},
                "columns": list(df.columns),
                "sample_data": df.head(5).to_dict('records'),
                "data_analysis": analysis,
                "has_formulas": self._has_formulas(worksheet),
                "formatted_cells": self._analyze_cell_formatting(worksheet)
            }
            
        except Exception as e:
            logger.warning(f"Sheet processing failed for {worksheet.title}: {e}")
            return {"name": worksheet.title, "error": str(e)}
    
    async def _detect_header_row(self, df: pd.DataFrame) -> Optional[int]:
        """Detect which row contains headers"""
        for i in range(min(5, len(df))):  # Check first 5 rows
            row = df.iloc[i]
            # Headers are usually strings and have fewer duplicates
            if row.dtype == 'object' and len(set(row.dropna())) > len(row) * 0.8:
                return i
        return None
    
    async def _analyze_tabular_data(self, df: pd.DataFrame, sheet_name: str) -> Dict[str, Any]:
        """Analyze tabular data for patterns and insights"""
        try:
            analysis = {
                "data_types": {},
                "missing_data": {},
                "unique_values": {},
                "potential_keys": [],
                "relationships": []
            }
            
            for col in df.columns:
                col_data = df[col].dropna()
                
                # Data type analysis
                analysis["data_types"][str(col)] = str(col_data.dtype)
                
                # Missing data
                missing_pct = (len(df) - len(col_data)) / len(df) * 100
                analysis["missing_data"][str(col)] = round(missing_pct, 2)
                
                # Unique values
                unique_count = len(col_data.unique())
                analysis["unique_values"][str(col)] = unique_count
                
                # Potential primary keys (unique values = total rows)
                if unique_count == len(col_data) and len(col_data) > 1:
                    analysis["potential_keys"].append(str(col))
            
            # Detect potential relationships between columns
            analysis["relationships"] = await self._detect_column_relationships(df)
            
            return analysis
            
        except Exception as e:
            logger.warning(f"Tabular data analysis failed: {e}")
            return {}
    
    async def _detect_column_relationships(self, df: pd.DataFrame) -> List[Dict[str, Any]]:
        """Detect relationships between columns"""
        relationships = []
        
        try:
            for i, col1 in enumerate(df.columns):
                for col2 in df.columns[i+1:]:
                    # Check for foreign key relationships
                    col1_values = set(df[col1].dropna().astype(str))
                    col2_values = set(df[col2].dropna().astype(str))
                    
                    # If one column's values are subset of another
                    if col1_values.issubset(col2_values) or col2_values.issubset(col1_values):
                        relationships.append({
                            "type": "potential_foreign_key",
                            "columns": [str(col1), str(col2)],
                            "confidence": 0.8
                        })
                    
                    # Check for correlation (numeric columns)
                    try:
                        if df[col1].dtype in ['int64', 'float64'] and df[col2].dtype in ['int64', 'float64']:
                            correlation = df[col1].corr(df[col2])
                            if abs(correlation) > 0.7:
                                relationships.append({
                                    "type": "correlation",
                                    "columns": [str(col1), str(col2)],
                                    "correlation": round(correlation, 3),
                                    "confidence": min(abs(correlation), 1.0)
                                })
                    except Exception:
                        pass
            
        except Exception as e:
            logger.warning(f"Relationship detection failed: {e}")
        
        return relationships
    
    def _has_formulas(self, worksheet) -> bool:
        """Check if worksheet contains formulas"""
        for row in worksheet.iter_rows():
            for cell in row:
                if cell.data_type == 'f':  # Formula
                    return True
        return False
    
    def _analyze_cell_formatting(self, worksheet) -> Dict[str, Any]:
        """Analyze cell formatting patterns"""
        formatting_info = {
            "has_merged_cells": len(worksheet.merged_cells.ranges) > 0,
            "has_conditional_formatting": len(worksheet.conditional_formatting) > 0,
            "formatted_cells": 0,
            "color_coded_cells": 0
        }
        
        for row in worksheet.iter_rows():
            for cell in row:
                if cell.fill.start_color.index != 0:  # Has background color
                    formatting_info["color_coded_cells"] += 1
                if cell.font.bold or cell.font.italic or cell.alignment.horizontal != 'general':
                    formatting_info["formatted_cells"] += 1
        
        return formatting_info
    
    async def _analyze_excel_chart(self, chart, sheet_name: str) -> Dict[str, Any]:
        """Analyze Excel chart properties"""
        try:
            return {
                "chart_type": type(chart).__name__,
                "title": getattr(chart, 'title', {}).get('tx', 'Untitled'),
                "sheet": sheet_name,
                "series_count": len(getattr(chart, 'series', [])),
                "has_legend": hasattr(chart, 'legend') and chart.legend is not None
            }
        except Exception as e:
            logger.warning(f"Chart analysis failed: {e}")
            return {"error": str(e)}
    
    async def _detect_pivot_tables(self, worksheet) -> List[Dict[str, Any]]:
        """Detect pivot tables in worksheet"""
        # This is a simplified detection - real pivot table detection would be more complex
        pivot_tables = []
        
        # Look for pivot table indicators (merged cells, specific formatting patterns)
        if len(worksheet.merged_cells.ranges) > 5:
            pivot_tables.append({
                "sheet": worksheet.title,
                "type": "potential_pivot_table",
                "confidence": 0.6
            })
        
        return pivot_tables
    
    async def _process_powerpoint(self, file_path: str, document_id: str) -> Dict[str, Any]:
        """Process PowerPoint presentations"""
        try:
            prs = Presentation(file_path)
            result = {
                "slides": [],
                "metadata": {
                    "total_slides": len(prs.slides),
                    "slide_layouts": []
                }
            }
            
            for i, slide in enumerate(prs.slides):
                slide_data = await self._process_powerpoint_slide(slide, i, document_id)
                result["slides"].append(slide_data)
                
                # Track layout usage
                layout_name = slide.slide_layout.name
                if layout_name not in result["metadata"]["slide_layouts"]:
                    result["metadata"]["slide_layouts"].append(layout_name)
            
            return result
            
        except Exception as e:
            logger.error(f"PowerPoint processing failed: {e}")
            raise DocumentProcessingError(f"PowerPoint processing failed: {e}")
    
    async def _process_powerpoint_slide(self, slide, slide_num: int, document_id: str) -> Dict[str, Any]:
        """Process individual PowerPoint slide"""
        try:
            slide_data = {
                "slide_number": slide_num + 1,
                "layout": slide.slide_layout.name,
                "text_content": [],
                "images": [],
                "charts": [],
                "tables": []
            }
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_data["text_content"].append({
                        "shape_type": str(shape.shape_type),
                        "text": shape.text.strip(),
                        "font_info": self._extract_font_info(shape) if hasattr(shape, 'text_frame') else {}
                    })
                
                # Extract images
                if shape.shape_type == 13:  # Picture
                    image_data = await self._extract_powerpoint_image(shape, slide_num, document_id)
                    if image_data:
                        slide_data["images"].append(image_data)
                
                # Extract charts
                if hasattr(shape, 'chart') and shape.chart:
                    chart_data = await self._analyze_powerpoint_chart(shape.chart, slide_num)
                    slide_data["charts"].append(chart_data)
                
                # Extract tables
                if hasattr(shape, 'table') and shape.table:
                    table_data = await self._extract_powerpoint_table(shape.table, slide_num)
                    slide_data["tables"].append(table_data)
            
            return slide_data
            
        except Exception as e:
            logger.warning(f"Slide {slide_num} processing failed: {e}")
            return {"slide_number": slide_num + 1, "error": str(e)}
    
    def _extract_font_info(self, shape) -> Dict[str, Any]:
        """Extract font information from text shape"""
        try:
            if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                para = shape.text_frame.paragraphs[0]
                if para.runs:
                    font = para.runs[0].font
                    return {
                        "name": font.name,
                        "size": font.size.pt if font.size else None,
                        "bold": font.bold,
                        "italic": font.italic
                    }
        except Exception:
            pass
        return {}
    
    async def _extract_powerpoint_image(self, shape, slide_num: int, document_id: str) -> Optional[Dict[str, Any]]:
        """Extract image from PowerPoint shape"""
        try:
            # This would require more complex image extraction from PowerPoint
            # For now, return basic information
            return {
                "image_id": f"slide_{slide_num}_image",
                "slide": slide_num + 1,
                "shape_type": "picture",
                "position": {
                    "left": shape.left.pt if shape.left else 0,
                    "top": shape.top.pt if shape.top else 0,
                    "width": shape.width.pt if shape.width else 0,
                    "height": shape.height.pt if shape.height else 0
                }
            }
        except Exception as e:
            logger.warning(f"PowerPoint image extraction failed: {e}")
            return None
    
    async def _analyze_powerpoint_chart(self, chart, slide_num: int) -> Dict[str, Any]:
        """Analyze PowerPoint chart"""
        try:
            return {
                "chart_type": str(chart.chart_type),
                "slide": slide_num + 1,
                "has_title": hasattr(chart, 'chart_title') and chart.chart_title,
                "series_count": len(chart.series) if hasattr(chart, 'series') else 0
            }
        except Exception as e:
            logger.warning(f"PowerPoint chart analysis failed: {e}")
            return {"error": str(e)}
    
    async def _extract_powerpoint_table(self, table, slide_num: int) -> Dict[str, Any]:
        """Extract table from PowerPoint"""
        try:
            table_data = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip())
                table_data.append(row_data)
            
            return {
                "slide": slide_num + 1,
                "dimensions": {"rows": len(table.rows), "columns": len(table.columns)},
                "data": table_data,
                "first_row_header": True  # Assumption
            }
        except Exception as e:
            logger.warning(f"PowerPoint table extraction failed: {e}")
            return {"error": str(e)}
    
    async def _process_csv(self, file_path: str, document_id: str) -> Dict[str, Any]:
        """Process CSV files with advanced analysis"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'latin1', 'cp1252']
            df = None
            
            for encoding in encodings:
                try:
                    df = pd.read_csv(file_path, encoding=encoding)
                    break
                except UnicodeDecodeError:
                    continue
            
            if df is None:
                raise DocumentProcessingError("Could not read CSV with any supported encoding")
            
            # Analyze CSV structure and content
            analysis = await self._analyze_tabular_data(df, "csv_data")
            
            return {
                "format": "csv",
                "dimensions": {"rows": len(df), "columns": len(df.columns)},
                "columns": list(df.columns),
                "sample_data": df.head(10).to_dict('records'),
                "data_analysis": analysis,
                "encoding_used": encoding,
                "data_quality": await self._assess_data_quality(df)
            }
            
        except Exception as e:
            logger.error(f"CSV processing failed: {e}")
            raise DocumentProcessingError(f"CSV processing failed: {e}")
    
    async def _assess_data_quality(self, df: pd.DataFrame) -> Dict[str, Any]:
        """Assess data quality metrics"""
        try:
            quality_metrics = {
                "completeness": {},
                "consistency": {},
                "accuracy": {},
                "overall_score": 0.0
            }
            
            total_cells = len(df) * len(df.columns)
            missing_cells = df.isnull().sum().sum()
            
            # Completeness
            quality_metrics["completeness"]["missing_percentage"] = (missing_cells / total_cells) * 100
            quality_metrics["completeness"]["score"] = max(0, 100 - quality_metrics["completeness"]["missing_percentage"])
            
            # Consistency (check for data type consistency)
            consistency_issues = 0
            for col in df.columns:
                if df[col].dtype == 'object':
                    # Check for mixed types in object columns
                    sample_values = df[col].dropna().astype(str)
                    if len(sample_values) > 0:
                        # Simple heuristic: if column has both numeric and non-numeric strings
                        numeric_count = sum(1 for val in sample_values if val.replace('.', '').replace('-', '').isdigit())
                        if 0 < numeric_count < len(sample_values):
                            consistency_issues += 1
            
            consistency_score = max(0, 100 - (consistency_issues / len(df.columns)) * 100)
            quality_metrics["consistency"]["score"] = consistency_score
            quality_metrics["consistency"]["issues"] = consistency_issues
            
            # Overall score
            quality_metrics["overall_score"] = (
                quality_metrics["completeness"]["score"] * 0.4 +
                quality_metrics["consistency"]["score"] * 0.6
            )
            
            return quality_metrics
            
        except Exception as e:
            logger.warning(f"Data quality assessment failed: {e}")
            return {"error": str(e)}
    
    async def _process_json(self, file_path: str, document_id: str) -> Dict[str, Any]:
        """Process JSON files with structure analysis"""
        try:
            import json
            
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            # Analyze JSON structure
            structure = await self._analyze_json_structure(data)
            
            return {
                "format": "json",
                "structure": structure,
                "sample_data": self._get_json_sample(data),
                "schema": await self._infer_json_schema(data)
            }
            
        except Exception as e:
            logger.error(f"JSON processing failed: {e}")
            raise DocumentProcessingError(f"JSON processing failed: {e}")
    
    async def _analyze_json_structure(self, data: Any) -> Dict[str, Any]:
        """Analyze JSON structure recursively"""
        def analyze_recursive(obj, path="root"):
            if isinstance(obj, dict):
                return {
                    "type": "object",
                    "keys": list(obj.keys()),
                    "properties": {k: analyze_recursive(v, f"{path}.{k}") for k, v in obj.items()}
                }
            elif isinstance(obj, list):
                if obj:
                    # Analyze first few items to understand array structure
                    sample_analysis = [analyze_recursive(item, f"{path}[{i}]") for i, item in enumerate(obj[:3])]
                    return {
                        "type": "array",
                        "length": len(obj),
                        "item_types": sample_analysis
                    }
                else:
                    return {"type": "array", "length": 0, "item_types": []}
            else:
                return {
                    "type": type(obj).__name__,
                    "value": str(obj)[:100] + "..." if len(str(obj)) > 100 else str(obj)
                }
        
        return analyze_recursive(data)
    
    def _get_json_sample(self, data: Any, max_items: int = 5) -> Any:
        """Get sample of JSON data for preview"""
        if isinstance(data, dict):
            return {k: self._get_json_sample(v, max_items) for k, v in list(data.items())[:max_items]}
        elif isinstance(data, list):
            return [self._get_json_sample(item, max_items) for item in data[:max_items]]
        else:
            return data
    
    async def _infer_json_schema(self, data: Any) -> Dict[str, Any]:
        """Infer JSON schema from data"""
        # This is a simplified schema inference
        # In production, you might use libraries like genson
        def infer_type(obj):
            if isinstance(obj, dict):
                return {
                    "type": "object",
                    "properties": {k: infer_type(v) for k, v in obj.items()}
                }
            elif isinstance(obj, list):
                if obj:
                    return {"type": "array", "items": infer_type(obj[0])}
                else:
                    return {"type": "array", "items": {}}
            elif isinstance(obj, bool):
                return {"type": "boolean"}
            elif isinstance(obj, int):
                return {"type": "integer"}
            elif isinstance(obj, float):
                return {"type": "number"}
            elif isinstance(obj, str):
                return {"type": "string"}
            elif obj is None:
                return {"type": "null"}
            else:
                return {"type": "unknown"}
        
        return infer_type(data)
    
    async def _process_html(self, file_path: str, document_id: str) -> Dict[str, Any]:
        """Process HTML files"""
        try:
            from bs4 import BeautifulSoup
            
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            soup = BeautifulSoup(content, 'html.parser')
            
            return {
                "format": "html",
                "title": soup.title.string if soup.title else "",
                "text_content": soup.get_text(),
                "structure": {
                    "headings": [{"level": tag.name, "text": tag.get_text()} 
                               for tag in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6'])],
                    "links": [{"text": a.get_text(), "href": a.get('href')} 
                             for a in soup.find_all('a', href=True)],
                    "images": [{"alt": img.get('alt', ''), "src": img.get('src', '')} 
                              for img in soup.find_all('img')],
                    "tables": len(soup.find_all('table')),
                    "forms": len(soup.find_all('form'))
                }
            }
            
        except Exception as e:
            logger.error(f"HTML processing failed: {e}")
            raise DocumentProcessingError(f"HTML processing failed: {e}")
    
    async def _process_xml(self, file_path: str, document_id: str) -> Dict[str, Any]:
        """Process XML files"""
        try:
            import xml.etree.ElementTree as ET
            
            tree = ET.parse(file_path)
            root = tree.getroot()
            
            # Analyze XML structure
            structure = self._analyze_xml_element(root)
            
            return {
                "format": "xml",
                "root_element": root.tag,
                "structure": structure,
                "namespaces": self._extract_xml_namespaces(root)
            }
            
        except Exception as e:
            logger.error(f"XML processing failed: {e}")
            raise DocumentProcessingError(f"XML processing failed: {e}")
    
    def _analyze_xml_element(self, element, max_depth=3, current_depth=0) -> Dict[str, Any]:
        """Analyze XML element structure"""
        if current_depth > max_depth:
            return {"type": "truncated"}
        
        analysis = {
            "tag": element.tag,
            "attributes": dict(element.attrib),
            "text": element.text.strip() if element.text else None,
            "children": {}
        }
        
        # Analyze child elements
        child_counts = {}
        for child in element:
            child_tag = child.tag
            if child_tag not in child_counts:
                child_counts[child_tag] = 0
                analysis["children"][child_tag] = self._analyze_xml_element(child, max_depth, current_depth + 1)
            child_counts[child_tag] += 1
        
        # Add count information
        for tag, count in child_counts.items():
            if tag in analysis["children"]:
                analysis["children"][tag]["count"] = count
        
        return analysis
    
    def _extract_xml_namespaces(self, root) -> Dict[str, str]:
        """Extract XML namespaces"""
        namespaces = {}
        for elem in root.iter():
            if elem.tag.startswith('{'):
                namespace = elem.tag[1:elem.tag.find('}')]
                local_name = elem.tag[elem.tag.find('}') + 1:]
                namespaces[local_name] = namespace
        return namespaces
    
    async def _process_image(self, file_path: str, document_id: str) -> Dict[str, Any]:
        """Process standalone image files"""
        try:
            image = Image.open(file_path)
            analysis = await self._analyze_image(image, document_id)
            
            return {
                "format": "image",
                "file_type": image.format,
                "dimensions": {"width": image.width, "height": image.height},
                "mode": image.mode,
                "analysis": analysis
            }
            
        except Exception as e:
            logger.error(f"Image processing failed: {e}")
            raise DocumentProcessingError(f"Image processing failed: {e}")
    
    async def _process_word(self, file_path: str, document_id: str) -> Dict[str, Any]:
        """Process Word documents"""
        try:
            from docx import Document as DocxDocument
            import os
            doc = DocxDocument(file_path)
            result = {
                "format": "word",
                "paragraphs": [],
                "tables": [],
                "images": [],
                "styles": set(),
                "metadata": {
                    "file_name": os.path.basename(file_path),
                    "num_paragraphs": len(doc.paragraphs),
                    "num_tables": len(doc.tables),
                }
            }

            # Extract paragraphs
            for para in doc.paragraphs:
                result["paragraphs"].append({
                    "text": para.text,
                    "style": para.style.name if para.style else None,
                    "runs": [run.text for run in para.runs],
                })
                if para.style:
                    result["styles"].add(para.style.name)

            # Extract tables
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                result["tables"].append({
                    "rows": len(table.rows),
                    "columns": len(table.columns),
                    "data": table_data
                })

            # Extract images (inline shapes)
            # python-docx does not support image extraction directly, so we note their presence
            # If you use docx2txt or other libs, you can extract images
            if hasattr(doc, 'inline_shapes'):
                for shape in doc.inline_shapes:
                    result["images"].append({
                        "type": "inline_shape",
                        "width": shape.width.pt if hasattr(shape.width, 'pt') else None,
                        "height": shape.height.pt if hasattr(shape.height, 'pt') else None
                    })

            # Document core properties
            props = doc.core_properties
            result["metadata"].update({
                "author": props.author,
                "created": str(props.created) if props.created else None,
                "last_modified_by": props.last_modified_by,
                "modified": str(props.modified) if props.modified else None,
                "title": props.title,
                "subject": props.subject,
                "category": props.category,
                "comments": props.comments,
            })

            # Convert styles to list
            result["styles"] = list(result["styles"])

            return result
        except Exception as e:
            logger.error(f"Word processing failed: {e}")
            raise DocumentProcessingError(f"Word processing failed: {e}")
            